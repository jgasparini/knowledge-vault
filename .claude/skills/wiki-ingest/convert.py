import sys
import pathlib
import zipfile
from html.parser import HTMLParser


def _require(module: str, package: str, ext: str) -> None:
    try:
        __import__(module)
    except ImportError:
        raise SystemExit(
            f"Missing dependency for {ext} files.\n"
            f"Install it with: pip install {package}"
        )


_HTML_BLOCK_TAGS = {"p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6"}

_MAX_SAMPLE_ROWS = 20


class _HTMLTextExtractor(HTMLParser):
    """Minimal HTML-to-text: drop tags, break lines at block-level elements."""

    def __init__(self):
        super().__init__()
        self._chunks = []

    def handle_starttag(self, tag, attrs):
        if tag in _HTML_BLOCK_TAGS:
            self._chunks.append("\n")

    def handle_data(self, data):
        self._chunks.append(data)

    def get_text(self) -> str:
        lines = (line.strip() for line in "".join(self._chunks).splitlines())
        return "\n".join(line for line in lines if line)


def _html_to_text(html: str) -> str:
    extractor = _HTMLTextExtractor()
    extractor.feed(html)
    return extractor.get_text()


def _iter_block_items(doc):
    """Yield each Paragraph and Table in document order (python-docx FAQ recipe)."""
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield Table(child, doc)


def _table_to_markdown(table) -> str:
    rows = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0:
            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(rows)


def _extract_media(src: pathlib.Path, zip_prefix: str, media_dir: pathlib.Path) -> list:
    """Extract files under zip_prefix from src (a zip-based Office file) into media_dir.

    Returns sorted filenames extracted; empty list (and no directory created) if none found.
    """
    with zipfile.ZipFile(src) as zf:
        names = [n for n in zf.namelist() if n.startswith(zip_prefix) and not n.endswith("/")]
        if not names:
            return []
        media_dir.mkdir(parents=True, exist_ok=True)
        extracted = []
        for name in names:
            filename = pathlib.Path(name).name
            (media_dir / filename).write_bytes(zf.read(name))
            extracted.append(filename)
    return sorted(extracted)


def _embedded_images_section(media_dir: pathlib.Path, filenames: list) -> str:
    lines = [
        "## Embedded images",
        f"This file contains {len(filenames)} embedded image(s), extracted to "
        f"`{media_dir}/`:",
    ]
    lines.extend(f"- {name}" for name in filenames)
    lines.append("")
    lines.append(
        "Read each image directly with the Read tool before writing the source "
        "summary, and note any diagrams, charts, or figures they show."
    )
    return "\n".join(lines)


def _cell_to_str(value) -> str:
    if value is None:
        return ""
    return str(value).replace("\n", " ").strip()


def _sheet_to_markdown(sheet) -> str:
    rows_iter = sheet.iter_rows(values_only=True)
    try:
        header = next(rows_iter)
    except StopIteration:
        return "_Empty sheet_"

    header_cells = [_cell_to_str(v) for v in header]
    data_rows = list(rows_iter)
    sample = data_rows[:_MAX_SAMPLE_ROWS]

    if len(data_rows) > _MAX_SAMPLE_ROWS:
        note = (
            f"_{len(header_cells)} columns x {len(data_rows)} data rows "
            f"(showing first {_MAX_SAMPLE_ROWS}). See the raw file for the full dataset._"
        )
    else:
        note = f"_{len(header_cells)} columns x {len(data_rows)} data rows_"

    lines = [
        note,
        "",
        "| " + " | ".join(header_cells) + " |",
        "| " + " | ".join(["---"] * len(header_cells)) + " |",
    ]
    for row in sample:
        lines.append("| " + " | ".join(_cell_to_str(v) for v in row) + " |")
    return "\n".join(lines)


def convert(src: pathlib.Path, media_dir: pathlib.Path = None) -> str:
    ext = src.suffix.lower()

    if ext == ".docx":
        _require("docx", "python-docx", ext)
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph
        doc = Document(src)
        lines = []
        for block in _iter_block_items(doc):
            if isinstance(block, Paragraph):
                if block.text.strip():
                    lines.append(block.text)
            elif isinstance(block, Table):
                lines.append(_table_to_markdown(block))
        if media_dir is not None:
            images = _extract_media(src, "word/media/", media_dir)
            if images:
                lines.append(_embedded_images_section(media_dir, images))
        return "\n\n".join(lines)

    if ext in (".pptx", ".ppt"):
        _require("pptx", "python-pptx", ext)
        from pptx import Presentation
        try:
            prs = Presentation(src)
        except Exception as e:
            if ext == ".ppt":
                raise SystemExit(
                    "python-pptx cannot open legacy .ppt files (binary format).\n"
                    "Open the file in PowerPoint or LibreOffice and export it as .pptx, "
                    "then ingest the .pptx version."
                ) from e
            raise
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        lines.append(text)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"_Notes: {notes}_")
        if media_dir is not None:
            images = _extract_media(src, "ppt/media/", media_dir)
            if images:
                lines.append(_embedded_images_section(media_dir, images))
        return "\n\n".join(lines)

    if ext == ".xlsx":
        _require("openpyxl", "openpyxl", ext)
        import openpyxl
        wb = openpyxl.load_workbook(src, data_only=True, read_only=True)
        parts = []
        for name in wb.sheetnames:
            parts.append(f"## Sheet: {name}")
            parts.append(_sheet_to_markdown(wb[name]))
        return "\n\n".join(parts)

    if ext == ".xls":
        raise SystemExit(
            "openpyxl cannot open legacy .xls files (binary format).\n"
            "Open the file in Excel or LibreOffice and re-export as .xlsx, "
            "then ingest the .xlsx version."
        )

    if ext == ".eml":
        import email as _email
        msg = _email.message_from_bytes(src.read_bytes())
        parts = [
            f"From: {msg.get('From', '')}",
            f"To: {msg.get('To', '')}",
            f"Subject: {msg.get('Subject', '')}",
            f"Date: {msg.get('Date', '')}",
            "",
        ]
        body = ""
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    body = payload.decode(charset, errors="replace")
                    break
        if not body:
            for part in msg.walk():
                if part.get_content_type() == "text/html":
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or "utf-8"
                        body = _html_to_text(payload.decode(charset, errors="replace"))
                        break
        if not body:
            raise SystemExit(
                f"No text/plain or text/html body found in {src}.\n"
                "Refusing to write a body-less file."
            )
        parts.append(body)
        return "\n".join(parts)

    if ext == ".msg":
        _require("extract_msg", "extract-msg", ext)
        import extract_msg
        with extract_msg.openMsg(src) as m:
            return "\n".join([
                f"From: {getattr(m, 'sender', '')}",
                f"To: {getattr(m, 'to', '')}",
                f"Subject: {getattr(m, 'subject', '')}",
                f"Date: {getattr(m, 'date', '')}",
                "",
                m.body or "",
            ])

    raise ValueError(f"Unsupported format: {ext}")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python convert.py <source_file>", file=sys.stderr)
        sys.exit(1)
    src = pathlib.Path(sys.argv[1])
    out = src.with_suffix(".md")
    if out.exists():
        print(f"Refusing to overwrite existing file: {out}", file=sys.stderr)
        sys.exit(1)
    media_dir = out.parent / f"{out.stem}_media"
    out.write_text(convert(src, media_dir=media_dir), encoding="utf-8")
    print(out)
