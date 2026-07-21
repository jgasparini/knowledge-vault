import pathlib
import subprocess
import sys
import pytest

from convert import convert


# ---------------------------------------------------------------------------
# .docx
# ---------------------------------------------------------------------------

def test_docx_extracts_paragraphs(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("Hello from Word")
    doc.add_paragraph("Second paragraph")
    src = tmp_path / "sample.docx"
    doc.save(src)

    result = convert(src)
    assert "Hello from Word" in result
    assert "Second paragraph" in result


def test_docx_skips_empty_paragraphs(tmp_path):
    from docx import Document

    doc = Document()
    doc.add_paragraph("Real text")
    doc.add_paragraph("")
    doc.add_paragraph("More text")
    src = tmp_path / "gaps.docx"
    doc.save(src)

    result = convert(src)
    assert "Real text" in result
    assert "More text" in result
    assert "\n\n\n" not in result  # no double-blank lines from empty paras


def test_docx_extracts_table_in_document_order(tmp_path):
    from docx import Document

    doc = Document()
    doc.add_paragraph("Before the table")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Row1 A"
    table.cell(1, 1).text = "Row1 B"
    doc.add_paragraph("After the table")
    src = tmp_path / "with_table.docx"
    doc.save(src)

    result = convert(src)
    assert "Header A" in result
    assert "Row1 B" in result
    before_idx = result.index("Before the table")
    table_idx = result.index("Header A")
    after_idx = result.index("After the table")
    assert before_idx < table_idx < after_idx


# ---------------------------------------------------------------------------
# .pptx
# ---------------------------------------------------------------------------

def test_pptx_extracts_text_frames(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide One Title"
    src = tmp_path / "deck.pptx"
    prs.save(src)

    result = convert(src)
    assert "Slide One Title" in result
    assert "Slide 1" in result


def test_pptx_includes_speaker_notes(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "These are speaker notes"
    src = tmp_path / "notes.pptx"
    prs.save(src)

    result = convert(src)
    assert "These are speaker notes" in result


# ---------------------------------------------------------------------------
# .xlsx
# ---------------------------------------------------------------------------

def test_xlsx_extracts_header_and_rows(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Score"])
    ws.append(["Alice", 90])
    ws.append(["Bob", 85])
    src = tmp_path / "sample.xlsx"
    wb.save(src)

    result = convert(src)
    assert "## Sheet: Data" in result
    assert "Name" in result and "Score" in result
    assert "Alice" in result and "90" in result
    assert "2 columns x 2 data rows" in result


def test_xlsx_multiple_sheets_both_present(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "First"
    ws1.append(["A"])
    ws1.append(["1"])
    ws2 = wb.create_sheet("Second")
    ws2.append(["B"])
    ws2.append(["2"])
    src = tmp_path / "multi.xlsx"
    wb.save(src)

    result = convert(src)
    assert "## Sheet: First" in result
    assert "## Sheet: Second" in result
    first_idx = result.index("## Sheet: First")
    second_idx = result.index("## Sheet: Second")
    assert first_idx < second_idx


def test_xlsx_empty_sheet_noted(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Empty"
    src = tmp_path / "empty.xlsx"
    wb.save(src)

    result = convert(src)
    assert "_Empty sheet_" in result


def test_xlsx_caps_sample_and_notes_truncation(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["Col"])
    for i in range(30):
        ws.append([f"row{i}"])
    src = tmp_path / "big.xlsx"
    wb.save(src)

    result = convert(src)
    assert "showing first 20" in result
    assert "row0" in result
    assert "row19" in result
    assert "row20" not in result
    assert "30 data rows" in result


def test_xls_raises_helpful_error(tmp_path):
    src = tmp_path / "legacy.xls"
    src.write_bytes(b"not a real xls file")
    with pytest.raises(SystemExit, match="re-export as .xlsx"):
        convert(src)


# ---------------------------------------------------------------------------
# .eml
# ---------------------------------------------------------------------------

def test_eml_extracts_body_and_headers(tmp_path):
    raw = (
        b"From: alice@example.com\r\n"
        b"To: bob@example.com\r\n"
        b"Subject: Test Email\r\n"
        b"Date: Mon, 01 Jan 2024 12:00:00 +0000\r\n"
        b"Content-Type: text/plain\r\n"
        b"\r\n"
        b"Hello this is the email body."
    )
    src = tmp_path / "message.eml"
    src.write_bytes(raw)

    result = convert(src)
    assert "alice@example.com" in result
    assert "Test Email" in result
    assert "Hello this is the email body" in result


def test_eml_multipart_returns_text_plain(tmp_path):
    raw = (
        b"From: sender@example.com\r\n"
        b"Subject: Multipart\r\n"
        b"Content-Type: multipart/alternative; boundary=\"boundary\"\r\n"
        b"\r\n"
        b"--boundary\r\n"
        b"Content-Type: text/plain\r\n"
        b"\r\n"
        b"Plain text part\r\n"
        b"--boundary\r\n"
        b"Content-Type: text/html\r\n"
        b"\r\n"
        b"<html>HTML part</html>\r\n"
        b"--boundary--\r\n"
    )
    src = tmp_path / "multi.eml"
    src.write_bytes(raw)

    result = convert(src)
    assert "Plain text part" in result
    assert "<html>" not in result


def test_eml_decodes_non_utf8_charset(tmp_path):
    body = "Café with naïve résumé".encode("windows-1252")
    raw = (
        b"From: sender@example.com\r\n"
        b"Subject: Accents\r\n"
        b'Content-Type: text/plain; charset="windows-1252"\r\n'
        b"\r\n" + body
    )
    src = tmp_path / "accents.eml"
    src.write_bytes(raw)

    result = convert(src)
    assert "Café with naïve résumé" in result
    assert "�" not in result  # no mojibake replacement characters


def test_eml_html_only_extracts_text(tmp_path):
    raw = (
        b"From: sender@example.com\r\n"
        b"Subject: HTML only\r\n"
        b'Content-Type: text/html; charset="utf-8"\r\n'
        b"\r\n"
        b"<html><body><p>Hello world</p><p>Second paragraph</p></body></html>"
    )
    src = tmp_path / "html_only.eml"
    src.write_bytes(raw)

    result = convert(src)
    assert "Hello world" in result
    assert "Second paragraph" in result
    assert "<p>" not in result


def test_eml_no_body_raises(tmp_path):
    raw = (
        b"From: sender@example.com\r\n"
        b"Subject: No body\r\n"
        b'Content-Type: image/png\r\n'
        b"\r\n"
        b"\x89PNG\r\n"
    )
    src = tmp_path / "no_body.eml"
    src.write_bytes(raw)

    with pytest.raises(SystemExit, match="No text/plain or text/html body"):
        convert(src)


# ---------------------------------------------------------------------------
# .msg
# ---------------------------------------------------------------------------

def test_msg_extracts_headers_and_body(tmp_path, monkeypatch):
    import extract_msg

    class FakeMsg:
        sender = "alice@example.com"
        to = "bob@example.com"
        subject = "Test MSG"
        date = "Mon, 01 Jan 2024 12:00:00 +0000"
        body = "Hello from the msg body"

        def __enter__(self):
            return self

        def __exit__(self, *exc_info):
            return False

    src = tmp_path / "message.msg"
    src.write_bytes(b"not a real msg file")
    monkeypatch.setattr(extract_msg, "openMsg", lambda path: FakeMsg())

    result = convert(src)
    assert "alice@example.com" in result
    assert "Test MSG" in result
    assert "Hello from the msg body" in result


# ---------------------------------------------------------------------------
# Error cases
# ---------------------------------------------------------------------------

def test_unsupported_format_raises_value_error(tmp_path):
    src = tmp_path / "mystery.xyz"
    src.write_text("content")
    with pytest.raises(ValueError, match="Unsupported format"):
        convert(src)


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def test_cli_writes_md_file_and_prints_path(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("CLI test paragraph")
    src = tmp_path / "cli_test.docx"
    doc.save(src)

    script = pathlib.Path(__file__).parent.parent / "convert.py"
    result = subprocess.run(
        [sys.executable, str(script), str(src)],
        capture_output=True, text=True
    )
    assert result.returncode == 0, result.stderr
    out_path = pathlib.Path(result.stdout.strip())
    assert out_path.exists()
    assert out_path.suffix == ".md"
    assert "CLI test paragraph" in out_path.read_text()
    assert out_path.parent == tmp_path


def test_cli_bad_usage_exits_nonzero():
    script = pathlib.Path(__file__).parent.parent / "convert.py"
    result = subprocess.run(
        [sys.executable, str(script)],
        capture_output=True, text=True
    )
    assert result.returncode != 0


def test_cli_refuses_to_overwrite_existing_output(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("New content")
    src = tmp_path / "collision.docx"
    doc.save(src)

    existing = tmp_path / "collision.md"
    existing.write_text("Pre-existing content that must survive")

    script = pathlib.Path(__file__).parent.parent / "convert.py"
    result = subprocess.run(
        [sys.executable, str(script), str(src)],
        capture_output=True, text=True
    )
    assert result.returncode != 0
    assert "collision.md" in result.stderr
    assert existing.read_text() == "Pre-existing content that must survive"
