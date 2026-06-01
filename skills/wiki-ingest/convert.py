import sys
import pathlib


def _require(module: str, package: str, ext: str) -> None:
    try:
        __import__(module)
    except ImportError:
        raise SystemExit(
            f"Missing dependency for {ext} files.\n"
            f"Install it with: pip install {package}"
        )


def convert(src: pathlib.Path) -> str:
    ext = src.suffix.lower()

    if ext == ".docx":
        _require("docx", "python-docx", ext)
        from docx import Document
        return "\n\n".join(p.text for p in Document(src).paragraphs if p.text.strip())

    if ext in (".pptx", ".ppt"):
        _require("pptx", "python-pptx", ext)
        from pptx import Presentation
        try:
            prs = Presentation(src)
        except Exception as e:
            if ext == ".ppt":
                raise SystemExit(
                    "python-pptx cannot open legacy .ppt files (binary format).\n"
                    "Open the file in PowerPoint or LibreOffice and export it as .pptx, "
                    "then ingest the .pptx version."
                ) from e
            raise
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        lines.append(text)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"_Notes: {notes}_")
        return "\n\n".join(lines)

    if ext == ".eml":
        import email as _email
        msg = _email.message_from_bytes(src.read_bytes())
        parts = [
            f"From: {msg.get('From', '')}",
            f"To: {msg.get('To', '')}",
            f"Subject: {msg.get('Subject', '')}",
            f"Date: {msg.get('Date', '')}",
            "",
        ]
        body = ""
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    body = payload.decode(errors="replace")
                    break
        if body:
            parts.append(body)
        return "\n".join(parts)

    if ext == ".msg":
        _require("extract_msg", "extract-msg", ext)
        import extract_msg
        with extract_msg.openMsg(src) as m:
            return "\n".join([
                f"From: {getattr(m, 'sender', '')}",
                f"To: {getattr(m, 'to', '')}",
                f"Subject: {getattr(m, 'subject', '')}",
                f"Date: {getattr(m, 'date', '')}",
                "",
                m.body or "",
            ])

    raise ValueError(f"Unsupported format: {ext}")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python convert.py <source_file>", file=sys.stderr)
        sys.exit(1)
    src = pathlib.Path(sys.argv[1])
    out = src.with_suffix(".md")
    out.write_text(convert(src), encoding="utf-8")
    print(out)
