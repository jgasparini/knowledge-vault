import pathlib
import subprocess
import sys
import pytest

from convert import convert


# ---------------------------------------------------------------------------
# .docx
# ---------------------------------------------------------------------------

def test_docx_extracts_paragraphs(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("Hello from Word")
    doc.add_paragraph("Second paragraph")
    src = tmp_path / "sample.docx"
    doc.save(src)

    result = convert(src)
    assert "Hello from Word" in result
    assert "Second paragraph" in result


def test_docx_skips_empty_paragraphs(tmp_path):
    from docx import Document

    doc = Document()
    doc.add_paragraph("Real text")
    doc.add_paragraph("")
    doc.add_paragraph("More text")
    src = tmp_path / "gaps.docx"
    doc.save(src)

    result = convert(src)
    assert "Real text" in result
    assert "More text" in result
    assert "\n\n\n" not in result  # no double-blank lines from empty paras


# ---------------------------------------------------------------------------
# .pptx
# ---------------------------------------------------------------------------

def test_pptx_extracts_text_frames(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide One Title"
    src = tmp_path / "deck.pptx"
    prs.save(src)

    result = convert(src)
    assert "Slide One Title" in result
    assert "Slide 1" in result


def test_pptx_includes_speaker_notes(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "These are speaker notes"
    src = tmp_path / "notes.pptx"
    prs.save(src)

    result = convert(src)
    assert "These are speaker notes" in result


# ---------------------------------------------------------------------------
# .eml
# ---------------------------------------------------------------------------

def test_eml_extracts_body_and_headers(tmp_path):
    raw = (
        b"From: alice@example.com\r\n"
        b"To: bob@example.com\r\n"
        b"Subject: Test Email\r\n"
        b"Date: Mon, 01 Jan 2024 12:00:00 +0000\r\n"
        b"Content-Type: text/plain\r\n"
        b"\r\n"
        b"Hello this is the email body."
    )
    src = tmp_path / "message.eml"
    src.write_bytes(raw)

    result = convert(src)
    assert "alice@example.com" in result
    assert "Test Email" in result
    assert "Hello this is the email body" in result


def test_eml_multipart_returns_text_plain(tmp_path):
    raw = (
        b"From: sender@example.com\r\n"
        b"Subject: Multipart\r\n"
        b"Content-Type: multipart/alternative; boundary=\"boundary\"\r\n"
        b"\r\n"
        b"--boundary\r\n"
        b"Content-Type: text/plain\r\n"
        b"\r\n"
        b"Plain text part\r\n"
        b"--boundary\r\n"
        b"Content-Type: text/html\r\n"
        b"\r\n"
        b"<html>HTML part</html>\r\n"
        b"--boundary--\r\n"
    )
    src = tmp_path / "multi.eml"
    src.write_bytes(raw)

    result = convert(src)
    assert "Plain text part" in result
    assert "<html>" not in result


# ---------------------------------------------------------------------------
# Error cases
# ---------------------------------------------------------------------------

def test_unsupported_format_raises_value_error(tmp_path):
    src = tmp_path / "mystery.xyz"
    src.write_text("content")
    with pytest.raises(ValueError, match="Unsupported format"):
        convert(src)


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def test_cli_writes_md_file_and_prints_path(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("CLI test paragraph")
    src = tmp_path / "cli_test.docx"
    doc.save(src)

    script = pathlib.Path(__file__).parent.parent / "convert.py"
    result = subprocess.run(
        [sys.executable, str(script), str(src)],
        capture_output=True, text=True
    )
    assert result.returncode == 0, result.stderr
    out_path = pathlib.Path(result.stdout.strip())
    assert out_path.exists()
    assert out_path.suffix == ".md"
    assert "CLI test paragraph" in out_path.read_text()
    assert out_path.parent == tmp_path


def test_cli_bad_usage_exits_nonzero():
    script = pathlib.Path(__file__).parent.parent / "convert.py"
    result = subprocess.run(
        [sys.executable, str(script)],
        capture_output=True, text=True
    )
    assert result.returncode != 0
